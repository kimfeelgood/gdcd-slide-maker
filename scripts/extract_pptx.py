#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Extract rough content + 박스/노출 markers from a minimally-filled .pptx so the
gdcd-slide-maker skill can design a finished interactive deck from it.

Usage:
    python3 extract_pptx.py <input.pptx> <workdir>

Writes:
    <workdir>/manifest.json   per-slide text, images, box buttons, reveal targets
    <workdir>/media/...       every embedded image (and video) exported

Marker convention (matched by number N):
    박스N / box N      -> an interactive neumorphic box button
    노출N / 노출 N /
    reveal N / show N  -> the image / video / URL revealed when box N is clicked

Markers are detected in a shape's name, its alt-text (descr), and its text — and
inside grouped shapes. The marker token is recorded but should be stripped from
any visible label by the skill.
"""
import sys, os, re, json, hashlib
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# English aliases use \b so PowerPoint's default shape names ("TextBox 1") don't
# false-match "box 1"; Korean markers (박스/노출) are the primary, reliable tags.
BOX_RE    = re.compile(r'(?:박스|버튼|\bbox)\s*0*(\d+)', re.IGNORECASE)
REVEAL_RE = re.compile(r'(?:노출|\breveal|\bshow)\s*0*(\d+)', re.IGNORECASE)
URL_RE    = re.compile(r'(https?://[^\s)>\]]+)', re.IGNORECASE)
VIDEO_EXT = ('.mp4', '.webm', '.mov', '.m4v', '.ogg')


def name_descr(shape):
    """cNvPr 'name' and 'descr' (alt text) for any shape."""
    name = shape.name or ''
    descr = ''
    try:
        cNvPr = shape._element.find('.//' + qn('p:cNvPr'))
        if cNvPr is None:
            cNvPr = shape._element.find('.//' + qn('a:cNvPr'))
        if cNvPr is not None:
            descr = cNvPr.get('descr', '') or ''
            if not name:
                name = cNvPr.get('name', '') or ''
    except Exception:
        pass
    return name, descr


def shape_text(shape):
    if not shape.has_text_frame:
        return ''
    return '\n'.join(p.text for p in shape.text_frame.paragraphs if p.text).strip()


def pos(shape):
    try:
        return [round(Emu(shape.left).inches, 2), round(Emu(shape.top).inches, 2),
                round(Emu(shape.width).inches, 2), round(Emu(shape.height).inches, 2)]
    except Exception:
        return None


def iter_leaves(shapes):
    """Yield leaf shapes, descending into groups."""
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_leaves(sh.shapes)
        else:
            yield sh


def detect(*texts):
    """Return (box_n, reveal_n) found across the given strings, or (None, None)."""
    blob = '  '.join(t for t in texts if t)
    b = BOX_RE.search(blob)
    r = REVEAL_RE.search(blob)
    return (int(b.group(1)) if b else None,
            int(r.group(1)) if r else None)


def strip_markers(text):
    if not text:
        return ''
    text = BOX_RE.sub('', text)
    text = REVEAL_RE.sub('', text)
    return text.strip(' \t\n·:-')


def save_image(shape, media_dir, tag, registry):
    """Save an image, de-duplicating identical blobs (authors often stack two
    copies of the same picture). Returns the media path, or None."""
    try:
        img = shape.image
    except Exception:
        return None
    h = hashlib.md5(img.blob).hexdigest()
    if h in registry:
        return registry[h]                     # identical image already saved
    ext = img.ext or 'png'
    fn = f'{tag}.{ext}'
    with open(os.path.join(media_dir, fn), 'wb') as f:
        f.write(img.blob)
    registry[h] = f'media/{fn}'
    return registry[h]


def slide_fingerprint(rec, path2md5):
    """A position-independent content hash of one slide. Stays stable when slides
    are reordered or new ones are inserted, so already-built slides can be matched
    and preserved. Excludes geometry (pos) — only content + markers matter."""
    parts = sorted(t['text'].strip() for t in rec['texts'])
    parts += sorted('img:' + path2md5.get(im.get('file'), im.get('file', '')) for im in rec['images'])
    for b in sorted(rec['boxes'], key=lambda e: e['n']):
        parts.append(f"box{b['n']}:{b.get('label') or ''}")
    for r in sorted(rec['reveals'], key=lambda e: e['n']):
        key = r.get('url') or r.get('text') or path2md5.get(r.get('file'), r.get('file', ''))
        parts.append(f"rev{r['n']}:{r.get('type')}:{key}")
    return hashlib.md5('\n'.join(parts).encode('utf-8')).hexdigest()[:12]


def _center(p):
    return (p[0] + p[2] / 2.0, p[1] + p[3] / 2.0)


def resolve_reveals(rec):
    """A 노출N marker is often a small label placed ON TOP of an image rather than
    the image itself being named. Bind each such empty-text reveal to the image it
    overlaps (or is nearest to), then drop redundant labels and dedupe."""
    imgs = rec['images']
    used = set()
    out = []
    for rv in rec['reveals']:
        if rv.get('type') == 'text' and not rv.get('text') and rv.get('pos'):
            x, y = rv['pos'][0], rv['pos'][1]
            best, bestd = None, 1e18
            for k, im in enumerate(imgs):
                ip = im.get('pos')
                if k in used or not ip:
                    continue
                inside = (ip[0] - 0.1 <= x <= ip[0] + ip[2] + 0.1 and
                          ip[1] - 0.1 <= y <= ip[1] + ip[3] + 0.1)
                if inside:
                    best, bestd = k, -1.0
                    break
                ic = _center(ip)
                d = (ic[0] - x) ** 2 + (ic[1] - y) ** 2
                if d < bestd:
                    bestd, best = d, k
            if best is not None:
                used.add(best)
                out.append({'n': rv['n'], 'type': 'image',
                            'file': imgs[best]['file'], 'pos': imgs[best]['pos']})
                continue
        out.append(rv)
    # drop empty-text labels whose number already resolved to a real asset
    real = {rv['n'] for rv in out
            if rv.get('type') in ('image', 'video', 'url')
            or (rv.get('type') == 'text' and rv.get('text'))}
    out = [rv for rv in out
           if not (rv.get('type') == 'text' and not rv.get('text') and rv['n'] in real)]
    # dedupe
    seen, dd = set(), []
    for rv in out:
        key = (rv['n'], rv.get('file'), rv.get('url'), rv.get('text'))
        if key in seen:
            continue
        seen.add(key)
        dd.append(rv)
    rec['reveals'] = dd
    rec['images'] = [im for k, im in enumerate(imgs) if k not in used]


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    src, workdir = sys.argv[1], sys.argv[2]
    media_dir = os.path.join(workdir, 'media')
    os.makedirs(media_dir, exist_ok=True)

    prs = Presentation(src)
    sw = round(Emu(prs.slide_width).inches, 2)
    sh_ = round(Emu(prs.slide_height).inches, 2)
    out = {'source': os.path.abspath(src),
           'slide_size_in': [sw, sh_],
           'slide_count': len(prs.slides),
           'slides': []}

    for si, slide in enumerate(prs.slides, 1):
        rec = {'index': si, 'texts': [], 'images': [], 'boxes': [], 'reveals': []}
        registry = {}                 # md5 -> media path, per slide (dedupe copies)
        added_files = set()
        for k, shape in enumerate(iter_leaves(slide.shapes), 1):
            nm, descr = name_descr(shape)
            txt = shape_text(shape)
            box_n, rev_n = detect(nm, descr, txt)
            is_pic = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            url_m = URL_RE.search(txt) if txt else None
            p = pos(shape)

            # --- reveal targets ---------------------------------------------
            if rev_n is not None:
                entry = {'n': rev_n, 'pos': p}
                if url_m:
                    u = url_m.group(1)
                    entry['type'] = 'video' if u.lower().endswith(VIDEO_EXT) else 'url'
                    entry['url'] = u
                    label = strip_markers(URL_RE.sub('', txt))
                    if label:
                        entry['label'] = label
                elif is_pic:
                    f = save_image(shape, media_dir, f's{si}_reveal{rev_n}', registry)
                    entry['type'] = 'image'
                    entry['file'] = f
                elif txt:
                    # reveal marker on a text box without a url: keep the text
                    entry['type'] = 'text'
                    entry['text'] = strip_markers(txt)
                else:
                    entry['type'] = 'unknown'
                rec['reveals'].append(entry)
                # a shape can be BOTH a box and a reveal only if numbers differ;
                # if same shape carries box+reveal, treat reveal here, box below
                if box_n is None:
                    continue

            # --- box buttons ------------------------------------------------
            if box_n is not None:
                entry = {'n': box_n, 'pos': p,
                         'label': strip_markers(txt) or None}
                if is_pic:
                    entry['face_image'] = save_image(shape, media_dir, f's{si}_box{box_n}', registry)
                rec['boxes'].append(entry)
                continue

            # --- plain (unmarked) content -----------------------------------
            if is_pic:
                f = save_image(shape, media_dir, f's{si}_img{k}', registry)
                if f and f not in added_files:        # skip duplicate copies
                    added_files.add(f)
                    rec['images'].append({'name': nm, 'file': f, 'pos': p})
            elif txt:
                rec['texts'].append({'name': nm, 'text': txt, 'pos': p})

        resolve_reveals(rec)          # bind 노출 labels to their images, dedupe

        # order boxes / reveals by number for convenience
        rec['boxes'].sort(key=lambda e: e['n'])
        rec['reveals'].sort(key=lambda e: e['n'])
        rec['fp'] = slide_fingerprint(rec, {v: k for k, v in registry.items()})
        out['slides'].append(rec)

    with open(os.path.join(workdir, 'manifest.json'), 'w', encoding='utf-8') as f:
        json.dump(out, f, ensure_ascii=False, indent=2)

    # short console summary
    print(f'slides: {out["slide_count"]}  ->  {os.path.join(workdir, "manifest.json")}')
    for s in out['slides']:
        print(f'  slide {s["index"]} [fp {s["fp"]}]: texts={len(s["texts"])} images={len(s["images"])} '
              f'boxes={[b["n"] for b in s["boxes"]]} reveals={[r["n"] for r in s["reveals"]]}')


if __name__ == '__main__':
    main()
